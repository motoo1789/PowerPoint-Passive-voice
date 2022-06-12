import pptx
from pptx.exc import PackageNotFoundError
from JUMAN import juman_passive_single
import glob

target_path = "./pptx/**"
pptx_files = []
files = glob.glob(target_path,recursive=True)

def pptx_start():

    target_text = []

    for ispptx in files:
        if ".pptx" in ispptx:
            pptx_files.append(ispptx)

    for pptx_file in pptx_files:
        print(pptx_file)

    len(pptx_files)

    extracted_pptx_texts = []
    extracted_pptx_pathes = []
    extracted_pptx_slide_num = []

    for pptx_file in pptx_files:
        # パスワードがかかっているかどうか
        try:
            # Presentation パワポの情報
            prs = pptx.Presentation(pptx_file)
        
        except PackageNotFoundError:
            print(pptx_file)
            print("パスワードがかかっているので開けません")
        
        for slide_num, slide in enumerate(prs.slides, start=0):
            for slide_shape in slide.shapes:

                if slide_shape.has_text_frame: # shapeに含まれるテキストデータ抽出
                    extracted_pptx_pathes.append(pptx_file)
                    extracted_pptx_slide_num.append(slide_num)
                    extracted_pptx_texts.append(slide_shape.text)
                    target_text.append(juman_passive_single(slide_shape.text))

                    target_text = [item for item in target_text if item != ""]
                    addNote(target_text)
                
                if slide_shape.has_table: #　テーブル
                    for cell in shape.table.iter_cells():
                        extracted_pptx_pathes.append(pptx_file)
                        extracted_pptx_slide_num.append(slide_num)
                        extracted_pptx_texts.append(cell.text)


# ノートの編集
def addNote(target_texts):

    for pptx_file in pptx_files:
        # パスワードがかかっているかどうか
        try:
            # Presentation パワポの情報
            prs = pptx.Presentation(pptx_file)
        
        except PackageNotFoundError:
            print(pptx_file)
            print("パスワードがかかっているので開けません")
        
        for slide_num, slide in enumerate(prs.slides, start=0):
            if slide.has_notes_slide:
                for slide_shape in slide.shapes:
                    for text in target_texts :
                        if text == slide_shape.text:
                            print(slide.notes_slide.notes_text_frame.text)
                            notes_slide = slide.notes_slide
                            text_frame = notes_slide.notes_text_frame
                            print(text_frame.text)
                            text_frame.text = "受け身だーーーーー"

        prs.save(pptx_file)

