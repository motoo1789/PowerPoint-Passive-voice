import pptx
from pptx.exc import PackageNotFoundError


prs = pptx.Presentation("./pptx/sample.pptx")

for slide_num, slide in enumerate(prs.slides, start=1):
    print(slide.notes_slide.notes_text_frame.text)

